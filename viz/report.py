"""Build a PPTX report from a generic JSON spec (python-pptx).

spec = {
  "title": "...", "subtitle": "...",
  "sections": [
    { "title": "...", "bullets": ["..", ".."],
      "images": [{"dataUrl": "data:image/png;base64,...", "caption": "..."}],
      "table": {"header": ["..", ".."], "rows": [["..", ".."], ...]} }
  ]
}
Each section becomes: one slide with title + bullets + images (1: full, 2: side by side, 3-4: 2x2 grid),
and an extra slide for the table when present.
"""
import base64, io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.5)


def _decode(data_url):
    b64 = data_url.split(",", 1)[1]
    return io.BytesIO(base64.b64decode(b64))


def _title(slide, text, top=Inches(0.3)):
    tb = slide.shapes.add_textbox(MARGIN, top, SW - 2 * MARGIN, Inches(0.7))
    p = tb.text_frame.paragraphs[0]; p.text = text
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = RGBColor(0x1F, 0x23, 0x28)
    return tb


def _bullets(slide, items, left, top, width, height, size=13):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for k, it in enumerate(items):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.text = "• " + str(it); p.font.size = Pt(size); p.space_after = Pt(4)
    return tb


def _picture_fit(slide, stream, left, top, box_w, box_h, caption=None):
    pic = slide.shapes.add_picture(stream, left, top)
    ratio = min(box_w / pic.width, (box_h - (Inches(0.35) if caption else 0)) / pic.height)
    pic.width = int(pic.width * ratio); pic.height = int(pic.height * ratio)
    pic.left = int(left + (box_w - pic.width) / 2); pic.top = int(top)
    if caption:
        tb = slide.shapes.add_textbox(left, pic.top + pic.height + Inches(0.05), box_w, Inches(0.3))
        p = tb.text_frame.paragraphs[0]; p.text = caption; p.font.size = Pt(11); p.font.color.rgb = RGBColor(0x6A, 0x70, 0x7A)
        p.alignment = 2  # center
    return pic


def _table(slide, header, rows, left, top, width, height):
    nrows, ncols = len(rows) + 1, len(header)
    shape = slide.shapes.add_table(nrows, ncols, left, top, width, min(height, Inches(0.32) * nrows))
    tbl = shape.table
    for c, h in enumerate(header):
        cell = tbl.cell(0, c); cell.text = str(h)
        cell.text_frame.paragraphs[0].font.size = Pt(11); cell.text_frame.paragraphs[0].font.bold = True
    for r, row in enumerate(rows, start=1):
        for c, v in enumerate(row):
            cell = tbl.cell(r, c); cell.text = "" if v is None else str(v)
            cell.text_frame.paragraphs[0].font.size = Pt(10)
    return shape


def build_pptx(spec, out_path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    blank = prs.slide_layouts[6]

    # title slide
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(MARGIN, Inches(2.4), SW - 2 * MARGIN, Inches(1.2))
    p = tb.text_frame.paragraphs[0]; p.text = spec.get("title", "iRIC 計算結果レポート"); p.font.size = Pt(40); p.font.bold = True
    if spec.get("subtitle"):
        tb = s.shapes.add_textbox(MARGIN, Inches(3.7), SW - 2 * MARGIN, Inches(1.5))
        tf = tb.text_frame; tf.word_wrap = True
        for k, line in enumerate(str(spec["subtitle"]).split("\n")):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.text = line; p.font.size = Pt(16); p.font.color.rgb = RGBColor(0x6A, 0x70, 0x7A)

    for sec in spec.get("sections", []):
        s = prs.slides.add_slide(blank)
        _title(s, sec.get("title", ""))
        top = Inches(1.1)
        bullets = sec.get("bullets") or []
        images = sec.get("images") or []
        content_h = SH - top - MARGIN
        if bullets and images:
            # bullets on the left third, images on the right
            _bullets(s, bullets, MARGIN, top, Inches(3.8), content_h)
            img_left, img_w = MARGIN + Inches(4.0), SW - 2 * MARGIN - Inches(4.0)
        else:
            if bullets:
                _bullets(s, bullets, MARGIN, top, SW - 2 * MARGIN, content_h)
            img_left, img_w = MARGIN, SW - 2 * MARGIN
        n = len(images)
        if n == 1:
            _picture_fit(s, _decode(images[0]["dataUrl"]), img_left, top, img_w, content_h, images[0].get("caption"))
        elif n == 2:
            w = (img_w - Inches(0.2)) / 2
            for k, im in enumerate(images):
                _picture_fit(s, _decode(im["dataUrl"]), img_left + k * (w + Inches(0.2)), top, w, content_h, im.get("caption"))
        elif n >= 3:
            cols = 2; rows = (n + 1) // 2
            w = (img_w - Inches(0.2)) / cols; h = (content_h - Inches(0.2) * (rows - 1)) / rows
            for k, im in enumerate(images[:6]):
                r, c = divmod(k, cols)
                _picture_fit(s, _decode(im["dataUrl"]), img_left + c * (w + Inches(0.2)), top + r * (h + Inches(0.2)), w, h, im.get("caption"))
        table = sec.get("table")
        if table and table.get("rows"):
            s2 = prs.slides.add_slide(blank)
            _title(s2, sec.get("title", "") + (" – 表" if images or bullets else ""))
            _table(s2, table["header"], table["rows"], MARGIN, Inches(1.1), SW - 2 * MARGIN, SH - Inches(1.6))
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    import sys, json
    spec = json.load(open(sys.argv[1], encoding="utf-8"))
    print(build_pptx(spec, sys.argv[2] if len(sys.argv) > 2 else "report.pptx"))
